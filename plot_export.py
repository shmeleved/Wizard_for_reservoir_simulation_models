import time
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from PyQt5.QtGui import QFont
import seaborn as sns
from pptx import Presentation
from pptx.util import Cm
import ecl_data_io as eclio
from PyQt5 import QtCore, QtWidgets, QtWebEngineWidgets
import plotly.graph_objects as go
import data_generation as dg
import os
import math

class Widget(QtWidgets.QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)

        self.button = QtWidgets.QPushButton('Построить', self)
        self.button.setFont(QFont('Arial', 12))
        self.export_button = QtWidgets.QPushButton('Экспорт', self)
        self.export_button.setFont(QFont('Arial', 12))
        self.well_button = QtWidgets.QPushButton('Выгрузка по \nскважинам', self)
        self.well_button.setFont(QFont('Arial', 12))
        self.report_button = QtWidgets.QPushButton('Сформировать \nотчет', self)
        self.report_button.setFont(QFont('Arial', 12))
        self.clear_cache_button = QtWidgets.QPushButton('Очистить \nкэш', self)
        self.clear_cache_button.setFont(QFont('Arial', 12))
        self.browser = QtWebEngineWidgets.QWebEngineView(self)

        self.well_button.move(870, 20)
        self.well_button.setFixedHeight(70)
        self.well_button.setFixedWidth(120)

        self.report_button.move(730, 20)
        self.report_button.setFixedHeight(70)
        self.report_button.setFixedWidth(120)

        self.clear_cache_button.move(590, 20)
        self.clear_cache_button.setFixedHeight(70)
        self.clear_cache_button.setFixedWidth(120)

        self.property_combo = QtWidgets.QComboBox(self)

        self.property_label = QtWidgets.QLabel('Свойство:', self)
        self.property_label.move(15, -5)
        self.property_combo.setFixedHeight(25)
        self.property_combo.setFixedWidth(120)
        self.property_combo.move(15, 20)

        self.time_combo = QtWidgets.QComboBox(self)
        self.time_label = QtWidgets.QLabel('Шаг:', self)
        self.time_label.move(15, 40)
        self.time_combo.setFixedHeight(25)
        self.time_combo.setFixedWidth(120)
        self.time_combo.move(15, 65)

        self.fip_combo = QtWidgets.QComboBox(self)
        self.fip_label = QtWidgets.QLabel('Номер/Имя:', self)
        self.fip_label.move(150, 40)
        self.fip_combo.setFixedHeight(25)
        self.fip_combo.setFixedWidth(120)
        self.fip_combo.move(150, 65)

        self.export_combo = QtWidgets.QComboBox(self)
        self.export_label = QtWidgets.QLabel('Метод экспорта:', self)
        self.export_label.move(285, 40)
        self.export_combo.setFixedHeight(25)
        self.export_combo.setFixedWidth(120)
        self.export_combo.move(285, 65)

        self.method_combo = QtWidgets.QComboBox(self)
        self.method_label = QtWidgets.QLabel('Метод построения:', self)
        self.method_label.move(285, -5)
        self.method_combo.setFixedHeight(25)
        self.method_combo.setFixedWidth(120)
        self.method_combo.move(285, 20)

        self.region_combo = QtWidgets.QComboBox(self)
        self.region_label = QtWidgets.QLabel('Регион:', self)
        self.region_label.move(150, -5)
        self.region_combo.setFixedHeight(25)
        self.region_combo.setFixedWidth(120)
        self.region_combo.move(150, 20)

        self.property = dg.get_static_property()
        self.dynamic_property = dg.get_dynamic_property()

        self.property_combo.addItem('Hэфф')
        self.property_combo.addItem('Hэфф(газ)')
        self.property_combo.addItem('Hэфф(нефть)')
        self.property_combo.addItem('Hэфф(нефть/газ)')

        remove = ['INTEHEAD', 'LOGIHEAD', 'DOUBHEAD','TABDIMS ', 'TAB     ','TNAVHEAD',
                  'IGRP    ', 'IWEL    ', 'SWEL    ', 'XWEL    ', 'ZWEL    ', 'ICON    ',
                  'SCON    ', 'XCON    ', 'STARTSOL','ENDSOL  ']

        for property_num in self.property:
            if property_num not in remove:
                self.property_combo.addItem(str(property_num))

        for i in range(len(self.dynamic_property)):
            if self.dynamic_property[i] not in self.property and self.dynamic_property[i] not in remove:
                self.property_combo.addItem(str(self.dynamic_property[i]))

        for time_step in list(dg.ts_reading().keys()):
            self.time_combo.addItem(time_step)

        self.export_combo.addItem('Раздельно')
        self.export_combo.addItem('Все регионы')

        self.method_combo.addItem('Сумма')
        self.method_combo.addItem('Среднее')
        self.method_combo.addItem('Средн.взвеш.(V)')
        self.method_combo.addItem('Средн.взвеш.(PORV)')
        self.method_combo.addItem('Плотность')
        self.method_combo.addItem('Минимум')

        self.region_combo.addItem('-')
        self.region_combo.addItem('FIPNUM')
        self.region_combo.addItem('EQLNUM')
        self.region_combo.addItem('PVTNUM')
        self.region_combo.addItem('SATNUM')
        self.region_combo.addItem('ROCKNUM')

        vlayout = QtWidgets.QVBoxLayout(self)
        vlayout.addWidget(self.button, alignment=QtCore.Qt.AlignHCenter)
        vlayout.addWidget(self.export_button, alignment=QtCore.Qt.AlignHCenter)
        self.browser.setMaximumHeight(680)
        vlayout.addWidget(self.browser)

        self.button.clicked.connect(self.show_graph)
        self.export_button.clicked.connect(self.export_cps)
        self.clear_cache_button.clicked.connect(dg.clear_cache)
        self.region_combo.currentIndexChanged.connect(self.num_region)
        self.resize(1000, 800)

        self.df_ijk, self.df_cps = dg.read_ijk()

        self.x_depart = dg.init_reading('DX')[0]
        self.y_depart = dg.init_reading('DY')[0]

        self.report_button.clicked.connect(self.report)

    def report(self):
        self.report_label = QtWidgets.QLabel('Тип отчета', self)
        time_report = self.time_combo.currentText()
        region_report = self.fip_combo.currentText()

        class MakeReport(QtWidgets.QWidget):
            def __init__(self, parent=None):
                super().__init__(parent)

                self.report_label = QtWidgets.QLabel('Тип отчета', self)

                self.report_1_button = QtWidgets.QPushButton('Общий отчет', self)
                self.report_2_button = QtWidgets.QPushButton('Для целей ПЗ', self)

                vlayout = QtWidgets.QVBoxLayout(self)
                vlayout.addWidget(self.report_label, alignment=QtCore.Qt.AlignHCenter)
                vlayout.addWidget(self.report_1_button)
                vlayout.addWidget(self.report_2_button)

                self.resize(200, 100)

                self.report_1_button.clicked.connect(self.report_1)

            def report_1(self):
                tic = time.perf_counter()
                if not os.path.isdir(dg.path_reports_export):
                    os.mkdir(dg.path_reports_export)

                if not os.path.isdir(dg.path_reports_export + '/01.Отчет_общий/'):
                    os.mkdir(dg.path_reports_export + '/01.Отчет_общий/')

                if not os.path.isdir(dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/'):
                    os.mkdir(dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/')

                if not dg.check_cache('cached/'+'df_report_1' + '_' + dg.model_name):
                    df, dummy = dg.read_ijk()
                    df.reset_index(inplace=True)
                    df['FIPNUM'] = dg.init_reading('FIPNUM')

                    df['Heff'] = np.array(dg.init_reading('NTG')) * np.array(dg.init_reading('DZ'))
                    df['VOL'] = np.array(dg.init_reading('DX')) * np.array(dg.init_reading('DY')) * np.array(
                        dg.init_reading('DZ'))
                    df['NTG'] = dg.init_reading('NTG')
                    df['PORO'] = dg.init_reading('PORO')
                    df['DZ'] = dg.init_reading('DZ')
                    df['DEPTH'] = dg.init_reading('DEPTH')
                    df['PERMX'] = dg.init_reading('PERMX')

                    for item in eclio.lazy_read(
                            dg.path_to_model + "/" + 'RESULTS/' + dg.model_name + '/' + dg.model_name + dg.x_files_list()[0]):
                        if 'FGIP' in item.read_keyword():
                            df['FGIP_init'] = item.read_array().astype(float)
                            print(df['FGIP_init'])

                    df.to_pickle('cached/'+'df_report_1' + '_' + dg.model_name)
                else:
                    df = pd.read_pickle('cached/'+'df_report_1' + '_' + dg.model_name).drop(columns='index')

                if not dg.check_cache('cached/'+'df_report_1' + '_' + dg.model_name + time_report):
                    df['SGAS'] = dg.unrst_reading('SGAS', time_report)
                    df['SOIL'] = dg.unrst_reading('SOIL', time_report)
                    df['FOIP'] = dg.unrst_reading('FOIP', time_report)
                    df['VOIP'] = dg.unrst_reading('VOIP', time_report)
                    df['FGIP'] = dg.unrst_reading('FGIP', time_report)
                    df['PRESSURE'] = dg.unrst_reading('PRESSURE', time_report)
                    df.to_pickle('cached/'+'df_report_1' + '_' + dg.model_name + time_report)
                else:
                    df = pd.read_pickle('cached/'+'df_report_1' + '_' + dg.model_name + time_report)

                region_number_list = dg.get_fip_dict(dg.path_to_fip_xlx)[region_report]

                index = ['Регион',
                         'Дата выгрузки',
                         'Средняя глубина залегания, м',
                         'Запасы свободного газа, млрд.м3',
                         'Запасы конденсата, тыс.т',
                         'Запасы нефти, тыс.т',
                         'Эффективный объем (газ), млн.м3',
                         'Эффективный поровый объем (газ), млн.м3',
                         'Эффективный поровый газонасыщенный объем, млн.м3',
                         'Эффективный объем (нефть), млн.м3',
                         'Эффективный поровый объем (нефть), млн.м3',
                         'Эффективный поровый нефтенасыщенный объем, млн.м3',
                         'Средн. эффективная газонасыщенная толщина, м',
                         'Средн. эффективная нефтенасыщенная толщина, м',
                         'Площадь газоносности, млн.м2',
                         'Площадь нефтеносности, млн.м2',
                         'Средн. пористость (газ), д.ед',
                         'Средн. пористость (нефть), д.ед',
                         'Средн. проницаемость (газ), мДарси',
                         'Средн. проницаемость (нефть), мДарси',
                         'Средн. газонасыщенность, д.ед',
                         'Средн. нефтенасыщенность, д.ед',
                         'Средн. пластовое давление (газ), бар',
                         'Средн. пластовое давление (нефть), бар']

                df_all = pd.DataFrame(index=index)
                df_all['0'] = range(24)

                df = df.query('FIPNUM in @region_number_list')

                df_gas = df.query('SGAS != 0')
                df_oil = df.query('SOIL != 0')

                df_all.iloc[0, 0] = region_report
                df_all.iloc[1, 0] = time_report
                df_all.iloc[2, 0] = round(df['DEPTH'].mean(),2)
                df_all.iloc[3, 0] = round(df['FGIP'].sum()/10**9,2)
                df_all.iloc[4, 0] = round(df['VOIP'].sum()/10**3,2)
                df_all.iloc[5, 0] = round(df['FOIP'].sum()/10**3,2)

                df_all.iloc[6, 0] = round(sum(df_gas['NTG'] * df_gas['VOL'])/10**6,2)
                df_all.iloc[7, 0] = round(sum(df_gas['NTG'] * df_gas['VOL'] * df_gas['PORO'])/10**6,2)
                df_all.iloc[8, 0] = round(sum(df_gas['NTG'] * df_gas['VOL'] * df_gas['PORO'] * df_gas['SGAS'])/10**6,2)

                df_all.iloc[9, 0] = round(sum(df_oil['NTG'] * df_oil['VOL'])/10**6,2)
                df_all.iloc[10, 0] = round(sum(df_oil['NTG'] * df_oil['VOL'] * df_oil['PORO'])/10**6,2)
                df_all.iloc[11, 0] = round(sum(df_oil['NTG'] * df_oil['VOL'] * df_oil['PORO'] * df_oil['SOIL'])/10**6,2)

                df_all.iloc[12, 0] = round(df_gas.groupby(['i', 'j'])['Heff'].sum().mean(),2)
                df_all.iloc[13, 0] = round(df_oil.groupby(['i', 'j'])['Heff'].sum().mean(),2)

                try:
                    df_all.iloc[14, 0] = round(df_all.iloc[6, 0] / df_all.iloc[12, 0],2)
                except ZeroDivisionError:
                    df_all.iloc[14, 0] = 0

                try:
                    df_all.iloc[15, 0] = round(df_all.iloc[9, 0] / df_all.iloc[13, 0],2)
                except ZeroDivisionError:
                    df_all.iloc[15, 0] = 0

                try:
                    df_all.iloc[16, 0] = round(df_all.iloc[7, 0] / df_all.iloc[6, 0],3)
                except ZeroDivisionError:
                    df_all.iloc[16, 0] = 0

                try:
                    df_all.iloc[17, 0] = round(df_all.iloc[10, 0] / df_all.iloc[9, 0],3)
                except ZeroDivisionError:
                    df_all.iloc[17, 0] = 0

                df_all.iloc[18, 0] = round(df_gas['PERMX'].mean(),2)
                df_all.iloc[19, 0] = round(df_oil['PERMX'].mean(),2)

                try:
                    df_all.iloc[20, 0] = round(df_all.iloc[8, 0] / df_all.iloc[7, 0],3)
                except ZeroDivisionError:
                    df_all.iloc[20, 0] = 0

                try:
                    df_all.iloc[21, 0] = round(df_all.iloc[11, 0] / df_all.iloc[10, 0],3)
                except ZeroDivisionError:
                    df_all.iloc[21, 0] = 0

                df_all.iloc[22, 0] = round(df_gas['PRESSURE'].mean(),2)
                df_all.iloc[23, 0] = round(df_oil['PRESSURE'].mean(),2)

                df_all.fillna(0, inplace=True)

                x_label = ['Начальные\n'+dg.get_key(dg.ts_reading(), dg.x_files_list()[0]), 'Текущие\n'+time_report]
                values = [round(df['FGIP_init'].sum()/10**9,2), round(df_all.iloc[3, 0],2)]

                def autolabel(rects, labels=None, height_factor=1.01):
                    try:
                        for i, rect in enumerate(rects):
                            height = rect.get_height()
                            if labels is not None:
                                try:
                                    label = labels[i]
                                except (TypeError, KeyError):
                                    label = ' '
                            else:
                                label = '%d' % int(height)
                            ax.text(rect.get_x() + rect.get_width() / 2., height_factor * height,
                                    '{}'.format(label),
                                    ha='center', va='bottom', size=20, weight='bold')
                    except IndexError:
                        pass

                bars = sns.barplot(x_label, values)
                ax = plt.gca()
                autolabel(ax.patches, values, height_factor=0.5)
                plt.ylabel('Запасы газа, млрд.м3')
                path_pic_1 = dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/' + region_report +'_'+ time_report + '_RESERVES.PNG'
                plt.savefig(path_pic_1,transparent=True)
                plt.close()

                df_gas['SGAS'].plot.hist(bins=100, alpha=0.8)
                plt.title('Гистограмма газонасыщенности пласта ПК1')
                plt.ylabel('Кол-во')
                plt.xlabel('Газонасыщенность, д.ед')

                mean_gas = df_gas['SGAS'].mean()

                plt.axvline(mean_gas, color='k', linestyle='dashed', linewidth=1)
                min_ylim, max_ylim = plt.ylim()
                plt.text(mean_gas * 1.01, max_ylim * 0.9, 'Mean: {:.2f}'.format(mean_gas), weight='bold')
                path_pic_2 = dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/' + region_report +'_'+ time_report + '_SGAS.PNG'
                plt.savefig(path_pic_2,transparent=True)
                plt.close()

                df['PORO'].plot.hist(bins=100, alpha=0.8)
                plt.title('Гистограмма пористости пласта ПК1')
                plt.ylabel('Кол-во')
                plt.xlabel('Пористость, д.ед')

                mean_poro = df['PORO'].mean()

                plt.axvline(mean_poro, color='k', linestyle='dashed',linewidth=1)
                min_ylim, max_ylim = plt.ylim()
                plt.text(mean_poro * 1.01, max_ylim * 0.9,'Mean: {:.2f}'.format(mean_poro), weight='bold')
                path_pic_3 = dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/' + region_report +'_'+ time_report + '_PORO.PNG'
                plt.savefig(path_pic_3,transparent=True)

                colorscale_iso = [[0, "rgb(255, 255, 255)"], [1, 'rgb(129, 193, 129)']]  # isobar
                colorscale_heff = [[0, "rgb(255, 255, 255)"], [1, 'rgb(175, 104, 0)']]  # heffg
                colorscale_struct = [[0, "rgb(182, 255, 255)"], [0.5, "rgb(255, 253, 253)"],
                                     [1, 'rgb(255, 182, 182)']]  # structure
                colorscale_res = [[0, "rgb(255, 255, 255)"], [0.1, "rgb(182, 255, 255)"], [0.5, "rgb(252, 254, 185)"],
                                  [1, 'rgb(255, 146, 146)']]  # reserves

                def property_check(colorscale, property, property_rus, i, df):
                    if property != 'Heff':
                        fig = go.Figure()
                        fig.update_layout(margin=dict(l=0, r=0, t=0, b=0))
                        if property == 'PRESSURE':
                            df = df.groupby(['i', 'j'])[property].mean().reset_index()

                        elif property == 'DEPTH':
                            df = df.groupby(['i', 'j'])[property].min().reset_index()
                            df[property] *= -1
                        elif property == 'FGIP':
                            x_depart = dg.init_reading('DX')[0]
                            y_depart = dg.init_reading('DY')[0]
                            df = df.groupby(['i', 'j'])[property].sum()
                            df = df / (x_depart * y_depart)
                            df = df.reset_index()
                        fig.add_trace(
                            go.Contour(x=df['i'], y=df['j'], z=df[property],
                                       colorscale=colorscale, connectgaps=False,
                                       colorbar=dict(title=property_rus,
                                                     titleside='right',
                                                     titlefont=dict(size=14, family='Arial, sans-serif')),
                                       contours=dict(showlabels=True, labelfont=dict(size=12, color='black'))))

                    else:
                        fig = go.Figure()
                        fig.update_layout(margin=dict(l=0, r=0, t=0, b=0))
                        df = df.groupby(['i', 'j'])[property].sum().reset_index()
                        fig.add_trace(
                            go.Contour(x=df['i'], y=df['j'], z=df[property],
                                       colorscale=colorscale, connectgaps=False,
                                       colorbar=dict(title=property_rus,
                                                     titleside='right',
                                                     titlefont=dict(size=14, family='Arial, sans-serif')),
                                       contours=dict(showlabels=True, labelfont=dict(size=12, color='black'),
                                                     start=0, end=math.ceil(df[property].max()))))

                    fig.update_layout(
                        paper_bgcolor='rgba(0,0,0,0)',
                        plot_bgcolor='rgba(0,0,0,0)'
                    )
                    fig.update_yaxes(title='y', visible=False, showticklabels=False)
                    fig.update_xaxes(title='x', visible=False, showticklabels=False)
                    fig.write_image(dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/' + region_report + '_'+ time_report +'_MAP_'+ property + '.PNG')
                    plt.close('all')

                color_list = [colorscale_iso,colorscale_heff,colorscale_struct,colorscale_res]
                property_list = ['PRESSURE', 'Heff', 'DEPTH', 'FGIP']
                property_rus_list = ['Давление, бар', 'Эфф.газ.толщина, м', 'Глубина, м', 'Плотность запасов газа, тыс.м3/м2']

                for i in range(len(property_list)):
                    print('начали цикл построения')
                    property_check(color_list[i], property_list[i], property_rus_list[i], i, df_gas)

                currencies_ppt = Presentation('templates/REPORT_1_template.pptx')
                slide = currencies_ppt.slides[0]

                shapes = slide.shapes

                for shape in shapes:
                    if shape.shape_type == 13:
                        shapes.element.remove(shape.element)

                pic1 = path_pic_1
                pic2 = path_pic_2
                pic3 = path_pic_3
                map1 = dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/' + region_report + '_'+ time_report +'_MAP_'+ property_list[0] + '.PNG'
                map2 = dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/' + region_report + '_'+ time_report +'_MAP_'+ property_list[1] + '.PNG'
                map3 = dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/' + region_report + '_'+ time_report +'_MAP_'+ property_list[2] + '.PNG'
                map4 = dg.path_reports_export + '/01.Отчет_общий/' + 'PNG/' + region_report + '_'+ time_report +'_MAP_'+ property_list[3] + '.PNG'

                top_pic = slide.shapes.add_picture(pic1, Cm(8.7), Cm(1.1), width=Cm(8.5))
                mid_pic = slide.shapes.add_picture(pic2, Cm(8.7), Cm(6.8), width=Cm(8.5))
                bottom_pic = slide.shapes.add_picture(pic3, Cm(8.7), Cm(12.8), width=Cm(8.5))
                map1_pic = slide.shapes.add_picture(map1, Cm(16.6), Cm(2.5), width=Cm(8.7), height=Cm(7.5))
                map2_pic = slide.shapes.add_picture(map2, Cm(25.2), Cm(2.5), width=Cm(8.7), height=Cm(7.5))
                map3_pic = slide.shapes.add_picture(map3, Cm(25.2), Cm(11.2), width=Cm(8.7), height=Cm(7.5))
                map4_pic = slide.shapes.add_picture(map4, Cm(16.6), Cm(11.2), width=Cm(8.7), height=Cm(7.5))

                shapes = slide.shapes
                text_box_list = []
                auto_shape_list = []
                table_list = []
                place_holder_list = []
                for shape_idx in range(len(shapes)):
                    shape = shapes[shape_idx]
                    if shape.shape_type == 17:
                        text_box_list.append(shape_idx)
                    if shape.shape_type == 1:
                        auto_shape_list.append(shape_idx)
                    if shape.shape_type == 19:
                        table_list.append(shape_idx)
                    if shape.shape_type == 14:
                        place_holder_list.append(shape_idx)

                names = ['Карта изобар на ' + time_report,
                         'Карта плотности остаточных \nподвижных запасов газа на ' + time_report,
                         'Структурная карта на \n ' + time_report, 'Карта эфф.гн.толщин на ' + time_report]

                for i in range(len(text_box_list)):
                    paragraph = shapes[text_box_list[i]].text_frame.paragraphs[0]
                    paragraph.runs[0].text = names[i]

                paragraph = shapes[place_holder_list[1]].text_frame.paragraphs[0]
                paragraph.runs[0].text = 'Общий отчет по пласту '+ region_report +' на ' + time_report

                table = shapes[table_list[0]].table

                for i in range(len(df_all)):
                    cell = table.cell(i, 1)
                    paragraph = cell.text_frame.paragraphs[0]
                    run = paragraph.runs[0]
                    run.text = str(df_all.iloc[i,0])

                currencies_ppt.save(dg.path_reports_export + '/01.Отчет_общий/' + region_report + '_' + time_report + '_Общий_отчет.pptx')

        self.w3 = MakeReport()
        self.w3.show()

    def num_region(self):
        self.fip_combo.clear()
        self.fip_combo.addItem('ALL')

        if os.path.exists(dg.path_to_model + '/'+self.region_combo.currentText()+'.xlsx') is True:
            dg.path_to_fip_xlx = dg.path_to_model + '/'+self.region_combo.currentText()+'.xlsx'
        else:
            dg.path_to_fip_xlx = None

        if dg.path_to_fip_xlx == None and self.region_combo.currentText() != '-':
            for item in eclio.lazy_read(
                    dg.path_to_model + "/" + 'RESULTS/' + dg.model_name + '/' + dg.model_name + '.INIT'):
                if self.region_combo.currentText() in item.read_keyword():
                    for fip in list(set(item.read_array())):
                        self.fip_combo.addItem(str(fip))
        elif dg.path_to_fip_xlx == None and self.region_combo.currentText() == '-':
            self.fip_combo.clear()
            self.fip_combo.addItem('-регион не выбран-')
        else:
            for object in dg.get_fip_dict(dg.path_to_fip_xlx).keys():
                self.fip_combo.addItem(object)

    def show_graph(self):
        tic = time.perf_counter()
        self.df_ijk_plot = self.df_ijk.copy()
        self.df_ijk_plot.reset_index(drop=True, inplace=True)
        self.df_ijk_plot['region'] = dg.init_reading(self.region_combo.currentText())

        if self.property_combo.currentText() == 'Hэфф':
            dznet = dg.get_dznet()
            self.df_ijk_plot[self.property_combo.currentText()] = dznet

        elif self.property_combo.currentText() == 'Hэфф(газ)':
            dznet = dg.get_dznet()
            self.df_ijk_plot[self.property_combo.currentText()] = dznet
            self.df_ijk_plot['SGAS'] = dg.unrst_reading('SGAS', self.time_combo.currentText())
            self.df_ijk_plot = self.df_ijk_plot.query('SGAS != 0')
            self.df_ijk_plot[self.property_combo.currentText()] = self.df_ijk_plot[self.property_combo.currentText()] * self.df_ijk_plot['SGAS']

        elif self.property_combo.currentText() == 'Hэфф(нефть)':
            dznet = dg.get_dznet()
            self.df_ijk_plot[self.property_combo.currentText()] = dznet
            self.df_ijk_plot['SOIL'] = dg.unrst_reading('SOIL', self.time_combo.currentText())
            self.df_ijk_plot = self.df_ijk_plot.query('SOIL != 0')
            self.df_ijk_plot[self.property_combo.currentText()] = self.df_ijk_plot[self.property_combo.currentText()] * self.df_ijk_plot['SOIL']

        elif self.property_combo.currentText() == 'Hэфф(нефть/газ)':
            dznet = dg.get_dznet()
            self.df_ijk_plot[self.property_combo.currentText()] = dznet
            sgas = dg.unrst_reading('SGAS', self.time_combo.currentText())
            soil = dg.unrst_reading('SOIL', self.time_combo.currentText())
            shc = sgas + soil
            self.df_ijk_plot['SHC'] = shc
            self.df_ijk_plot = self.df_ijk_plot.query('SHC != 0')
            self.df_ijk_plot[self.property_combo.currentText()] = self.df_ijk_plot[self.property_combo.currentText()] * self.df_ijk_plot['SHC']

        elif self.property_combo.currentText() in self.dynamic_property:
            self.df_ijk_plot[self.property_combo.currentText()] = dg.unrst_reading(self.property_combo.currentText(), self.time_combo.currentText())
        else:
            self.df_ijk_plot[self.property_combo.currentText()] = dg.init_reading(self.property_combo.currentText())

        self.df_ijk_all = self.df_ijk_plot.copy()

        if dg.path_to_fip_xlx == None:
            if self.fip_combo.currentText() != 'ALL':
                fip = int(self.fip_combo.currentText())
                self.df_ijk_plot = self.df_ijk_plot.query('region == @fip')
            if self.method_combo.currentText() == 'Среднее':
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText()].mean().reset_index()
            elif self.method_combo.currentText() == 'Сумма':
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText()].sum().reset_index()
            elif self.method_combo.currentText() == 'Плотность':
                x_depart = dg.init_reading('DX')[0]
                y_depart = dg.init_reading('DY')[0]
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText()].sum()
                self.grid = self.grid / (self.x_depart * self.y_depart)
                self.grid = self.grid.reset_index()
            elif self.method_combo.currentText() == 'Средн.взвеш.(V)':
                self.df_ijk_plot['volume'] = dg.get_volume()
                self.df_ijk_plot[self.property_combo.currentText()] = self.df_ijk_plot[self.property_combo.currentText()] * self.df_ijk_plot['volume']
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText(),'volume'].sum()
                self.grid[self.property_combo.currentText()] = self.grid[self.property_combo.currentText()] / self.grid['volume']
                self.grid = self.grid.reset_index()
            elif self.method_combo.currentText() == 'Средн.взвеш.(PORV)':
                self.df_ijk_plot['volume'] = dg.init_reading('PORV')
                self.df_ijk_plot[self.property_combo.currentText()] = self.df_ijk_plot[self.property_combo.currentText()] * self.df_ijk_plot['volume']
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText(),'volume'].sum()
                self.grid[self.property_combo.currentText()] = self.grid[self.property_combo.currentText()] / self.grid['volume']
                self.grid = self.grid.reset_index()
            elif self.method_combo.currentText() == 'Минимум':
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText()].min().reset_index()
                self.grid[self.property_combo.currentText()] = self.grid[self.property_combo.currentText()]*(-1)
        else:
            if self.fip_combo.currentText() != 'ALL':
                fip = dg.get_fip_dict(dg.path_to_fip_xlx)[self.fip_combo.currentText()]
                self.df_ijk_plot = self.df_ijk_plot.query('region in @fip')
            if self.method_combo.currentText() == 'Среднее':
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText()].mean().reset_index()
            elif self.method_combo.currentText() == 'Сумма':
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText()].sum().reset_index()
            elif self.method_combo.currentText() == 'Плотность':
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText()].sum()
                self.grid = self.grid/(self.x_depart*self.y_depart)
                self.grid = self.grid.reset_index()
            elif self.method_combo.currentText() == 'Средн.взвеш.(V)':
                self.df_ijk_plot['volume'] = dg.get_volume()
                self.df_ijk_plot[self.property_combo.currentText()] = self.df_ijk_plot[self.property_combo.currentText()] * self.df_ijk_plot['volume']
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText(),'volume'].sum()
                self.grid[self.property_combo.currentText()] = self.grid[self.property_combo.currentText()] / self.grid['volume']
                self.grid = self.grid.reset_index()
            elif self.method_combo.currentText() == 'Средн.взвеш.(PORV)':
                self.df_ijk_plot['volume'] = dg.init_reading('PORV')
                self.df_ijk_plot[self.property_combo.currentText()] = self.df_ijk_plot[self.property_combo.currentText()] * self.df_ijk_plot['volume']
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText(),'volume'].sum()
                self.grid[self.property_combo.currentText()] = self.grid[self.property_combo.currentText()] / self.grid['volume']
                self.grid = self.grid.reset_index()
            elif self.method_combo.currentText() == 'Минимум':
                self.grid = self.df_ijk_plot.groupby(['i', 'j'])[self.property_combo.currentText()].min().reset_index()
                self.grid[self.property_combo.currentText()] = self.grid[self.property_combo.currentText()] * (-1)

        fig = go.Figure()
        config = dict({'scrollZoom': True})
        colorscale_iso = [[0, "rgb(255, 255, 255)"], [1, 'rgb(129, 193, 129)']] # isobar
        colorscale_heff = [[0, "rgb(255, 255, 255)"], [1, 'rgb(175, 104, 0)']] # heffg
        colorscale_struct = [[0, "rgb(182, 255, 255)"], [0.5, "rgb(255, 253, 253)"], [1, 'rgb(255, 182, 182)']] # structure
        colorscale_res = [[0, "rgb(255, 255, 255)"], [0.1, "rgb(182, 255, 255)"], [0.5, "rgb(252, 254, 185)"], [1, 'rgb(255, 146, 146)']] # reserves
        fig.update_layout(margin=dict(l=0, r=0, t=0, b=0))

        def property_check(colorscale):
            fig.add_trace(go.Contour(x=self.grid['i'], y=self.grid['j'], z=self.grid[self.property_combo.currentText()],
                                     colorscale=colorscale, connectgaps=False,
                                     colorbar=dict(title=self.property_combo.currentText(),
                                                   titleside='right',
                                                   titlefont=dict(size=14, family='Arial, sans-serif')),
                                     contours=dict(showlabels=True, labelfont=dict(size=12, color='black'))))

        if 'PRESSURE' in self.property_combo.currentText():
            property_check(colorscale_iso)
        elif 'Hэфф' in self.property_combo.currentText():
            property_check(colorscale_heff)
        elif 'DEPTH' in self.property_combo.currentText():
            property_check(colorscale_struct)
        elif 'FIPGAS' in self.property_combo.currentText():
            property_check(colorscale_res)
        else:
            property_check('viridis_r')

        fig.update_layout(
            paper_bgcolor='rgba(0,0,0,0)',
            plot_bgcolor='rgba(0,0,0,0)'
        )
        fig.update_yaxes(title='y', visible=False, showticklabels=False)
        fig.update_xaxes(title='x', visible=False, showticklabels=False)

        self.browser.setHtml(fig.to_html(include_plotlyjs='cdn', config=config))

    def export_cps(self):
        dg.export_cps(self.export_combo.currentText(),self.grid,self.property_combo.currentText(),
                      self.df_cps,self.time_combo.currentText(),self.fip_combo.currentText(),
                      self.df_ijk_all,self.region_combo.currentText(),self.method_combo.currentText())